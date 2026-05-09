from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


STYLE_PRESETS = {
    "business": {"bg": RGBColor(255, 255, 255), "title": RGBColor(28, 63, 170), "text": RGBColor(30, 30, 30)},
    "minimal": {"bg": RGBColor(245, 245, 245), "title": RGBColor(0, 0, 0), "text": RGBColor(50, 50, 50)},
    "creative": {"bg": RGBColor(35, 35, 70), "title": RGBColor(255, 200, 0), "text": RGBColor(240, 240, 240)},
    "education": {"bg": RGBColor(237, 248, 255), "title": RGBColor(0, 102, 153), "text": RGBColor(20, 20, 20)},
}


def build_ppt(slides: list[dict], style: str, output_path: str, title: str | None = None) -> None:
    preset = STYLE_PRESETS.get(style, STYLE_PRESETS["business"])

    prs = Presentation()
    if title:
        _add_title_slide(prs, title, preset)

    for slide in slides:
        _add_content_slide(prs, slide.get("title", "内容"), slide.get("bullets", []), preset)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _add_title_slide(prs: Presentation, title: str, preset: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_background(slide, preset)
    title_shape = slide.shapes.title
    title_shape.text = title
    run = title_shape.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = preset["title"]


def _add_content_slide(prs: Presentation, heading: str, bullets: list[str], preset: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_background(slide, preset)

    title_shape = slide.shapes.title
    title_shape.text = heading
    title_run = title_shape.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.color.rgb = preset["title"]

    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = str(bullet)
        p.level = 0
        if p.runs:
            p.runs[0].font.size = Pt(20)
            p.runs[0].font.color.rgb = preset["text"]


def _set_background(slide, preset: dict) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = preset["bg"]
